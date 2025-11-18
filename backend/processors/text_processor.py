# Extract text from documents   
"""
Text extraction from various document formats
"""
from pathlib import Path
from typing import List, Optional
from dataclasses import dataclass
from backend.utils.logger import app_logger as logger

try:
    import PyPDF2
    import pdfplumber
except ImportError:
    logger.warning("PDF libraries not installed")

try:
    from docx import Document
except ImportError:
    logger.warning("python-docx not installed")

try:
    from pptx import Presentation
except ImportError:
    logger.warning("python-pptx not installed")


@dataclass
class ProcessingResult:
    """Result container"""
    success: bool
    data: Optional[str] = None
    error: Optional[str] = None


class TextProcessor:
    """Extract text from various document formats"""
    
    def __init__(self, max_length: int = 50000):
        self.max_length = max_length
    
    def process(self, file_path: str) -> ProcessingResult:
        """Main processing entry point"""
        try:
            path = Path(file_path)
            suffix = path.suffix.lower()
            
            if suffix in ['.txt', '.md']:
                text = self._extract_from_text(file_path)
            elif suffix == '.pdf':
                text = self._extract_from_pdf(file_path)
            elif suffix in ['.doc', '.docx']:
                text = self._extract_from_word(file_path)
            elif suffix in ['.ppt', '.pptx']:
                text = self._extract_from_powerpoint(file_path)
            else:
                return ProcessingResult(False, error=f"Unsupported format: {suffix}")
            
            # Truncate if needed
            # if len(text) > self.max_length:
            #     text = text[:self.max_length] + "... (truncated)"
            
            return ProcessingResult(True, data=text)
            
        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}")
            return ProcessingResult(False, error=str(e))
    
    def _extract_from_text(self, file_path: str) -> str:
        """Extract from plain text files"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF"""
        text_parts = []
        
        try:
            # Try pdfplumber first (better for complex PDFs)
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
        except Exception as e:
            logger.warning(f"pdfplumber failed, trying PyPDF2: {e}")
            
            # Fallback to PyPDF2
            try:
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(page_text)
            except Exception as e2:
                logger.error(f"PyPDF2 also failed: {e2}")
                raise
        
        return "\n\n".join(text_parts)
    
    def _extract_from_word(self, file_path: str) -> str:
        """Extract text from Word documents"""
        doc = Document(file_path)
        
        text_parts = []
        
        # Extract paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells)
                text_parts.append(row_text)
        
        return "\n".join(text_parts)
    
    def _extract_from_powerpoint(self, file_path: str) -> str:
        """Extract text from PowerPoint"""
        prs = Presentation(file_path)
        
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n--- Slide {slide_num} ---")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        
        return "\n".join(text_parts)
    
    def extract_embedded_images(self, file_path: str) -> List[str]:
        """Extract paths to embedded images"""
        path = Path(file_path)
        suffix = path.suffix.lower()
        
        image_paths = []
        
        try:
            if suffix == '.pdf':
                image_paths = self._extract_images_from_pdf(file_path)
            elif suffix in ['.docx']:
                image_paths = self._extract_images_from_word(file_path)
            elif suffix in ['.pptx']:
                image_paths = self._extract_images_from_powerpoint(file_path)
        except Exception as e:
            logger.error(f"Error extracting images from {file_path}: {e}")
        
        return image_paths
    
    def _extract_images_from_pdf(self, file_path: str) -> List[str]:
        """Extract images from PDF"""
        import fitz  # PyMuPDF
        
        doc = fitz.open(file_path)
        image_paths = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            images = page.get_images()
            
            for img_index, img in enumerate(images):
                xref = img[0]
                base_image = doc.extract_image(xref)
                
                # Save image temporarily
                image_ext = base_image["ext"]
                image_path = Path(file_path).parent / f"{Path(file_path).stem}_page{page_num}_img{img_index}.{image_ext}"
                
                with open(image_path, "wb") as img_file:
                    img_file.write(base_image["image"])
                
                image_paths.append(str(image_path))
        
        return image_paths
    
    def _extract_images_from_word(self, file_path: str) -> List[str]:
        """Extract images from Word document"""
        doc = Document(file_path)
        image_paths = []
        
        # Extract inline images
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image_path = Path(file_path).parent / f"{Path(file_path).stem}_{rel.target_ref.split('/')[-1]}"
                
                with open(image_path, "wb") as img_file:
                    img_file.write(rel.target_part.blob)
                
                image_paths.append(str(image_path))
        
        return image_paths
    
    def _extract_images_from_powerpoint(self, file_path: str) -> List[str]:
        """Extract images from PowerPoint"""
        prs = Presentation(file_path)
        image_paths = []
        
        for slide_num, slide in enumerate(prs.slides):
            for shape_num, shape in enumerate(slide.shapes):
                if hasattr(shape, "image"):
                    image = shape.image
                    image_ext = image.ext
                    
                    image_path = Path(file_path).parent / f"{Path(file_path).stem}_slide{slide_num}_shape{shape_num}.{image_ext}"
                    
                    with open(image_path, "wb") as img_file:
                        img_file.write(image.blob)
                    
                    image_paths.append(str(image_path))
        
        return image_paths